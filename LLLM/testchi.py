import json
import io

import pdfplumber
from pptx import Presentation

from miyya import Chat

SYSTEM_TEST = """Sen test tuzuvchisan. Berilgan matn asosida o'zbek tilida test tuz.
Natijani FAQAT quyidagi JSON formatda qaytar, hech qanday izoh yozma:
[
  {
    "savol": "savol matni",
    "variantlar": ["A) ...", "B) ...", "C) ...", "D) ..."],
    "togri": "A"
  }
]"""


class Tuzuvchi:
    def __init__(self, soni: int = 10):
        self.soni = soni
        self._chat = Chat(sp=SYSTEM_TEST)

    def _pdfdan_matn(self, bayt: bytes) -> str:
        qismlar = []
        with pdfplumber.open(io.BytesIO(bayt)) as pdf:
            for sahifa in pdf.pages:
                m = sahifa.extract_text()
                if m:
                    qismlar.append(m)
        return '\n'.join(qismlar)

    def _pptxdan_matn(self, bayt: bytes) -> str:
        qismlar = []
        prs = Presentation(io.BytesIO(bayt))
        for slayd in prs.slides:
            for shakl in slayd.shapes:
                if shakl.has_text_frame:
                    for paragraf in shakl.text_frame.paragraphs:
                        m = paragraf.text.strip()
                        if m:
                            qismlar.append(m)
        return '\n'.join(qismlar)

    def _json_ajrat(self, javob: str) -> list:
        boshi = javob.find('[')
        ohiri = javob.rfind(']') + 1
        if boshi == -1 or ohiri == 0:
            return []
        return json.loads(javob[boshi:ohiri])

    def test_tuz(self, matn: str) -> list:
        prompt = f"{self.soni} ta savol tuz:\n\n{matn[:4000]}"
        javob = self._chat.sora(prompt)
        self._chat.tozala()
        return self._json_ajrat(javob)

    def fayldan_test(self, bayt: bytes, fayl_turi: str) -> list:
        if fayl_turi == 'pdf':
            matn = self._pdfdan_matn(bayt)
        elif fayl_turi == 'pptx':
            matn = self._pptxdan_matn(bayt)
        else:
            raise ValueError(f"Qo'llab-quvvatlanmagan tur: {fayl_turi}")

        if not matn.strip():
            raise ValueError("Fayldan matn chiqarib bo'lmadi")

        return self.test_tuz(matn)
